"""
PowerPoint Generator Module
Creates PowerPoint presentations based on template and content
"""

import io
from typing import Dict, List, Optional
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_SHAPE

from .chart_generator import ChartGenerator


class PresentationGenerator:
    """Generates PowerPoint presentations"""
    
    def __init__(self, template_data: Optional[Dict] = None):
        self.template_data = template_data or self._get_default_template()
        self.chart_generator = ChartGenerator()
        
    def _get_default_template(self) -> Dict:
        return {
            'colors': {
                'primary': '#6366F1',
                'secondary': '#8B5CF6',
                'accent': '#EC4899',
                'background': '#0F172A',
                'text_primary': '#F8FAFC',
                'text_secondary': '#94A3B8'
            },
            'fonts': {
                'title': {'name': 'Arial', 'size': 44, 'bold': True},
                'subtitle': {'name': 'Arial', 'size': 28, 'bold': False},
                'body': {'name': 'Arial', 'size': 18, 'bold': False}
            },
            'slide_size': {'width': 13.333, 'height': 7.5}
        }
    
    def generate_presentation(self, content: Dict, custom_settings: Optional[Dict] = None) -> io.BytesIO:
        """Generate a PowerPoint presentation from content"""
        
        if custom_settings:
            self._apply_custom_settings(custom_settings)
        
        prs = Presentation()
        prs.slide_width = Inches(self.template_data['slide_size']['width'])
        prs.slide_height = Inches(self.template_data['slide_size']['height'])
        
        slides = content.get('slides', [])
        for slide_content in slides:
            self._create_slide(prs, slide_content)
        
        output = io.BytesIO()
        prs.save(output)
        output.seek(0)
        
        return output
    
    def _apply_custom_settings(self, settings: Dict):
        if 'colors' in settings:
            self.template_data['colors'].update(settings['colors'])
        if 'fonts' in settings:
            for font_type, font_settings in settings['fonts'].items():
                if font_type in self.template_data['fonts']:
                    self.template_data['fonts'][font_type].update(font_settings)
    
    def _create_slide(self, prs: Presentation, slide_content: Dict):
        layout_type = slide_content.get('layout', 'content')
        
        blank_layout = prs.slide_layouts[6]
        slide = prs.slides.add_slide(blank_layout)
        
        self._add_background(slide)
        
        layout_creators = {
            'title': self._create_title_slide,
            'content': self._create_content_slide,
            'two_column': self._create_two_column_slide,
            'chart': self._create_chart_slide,
            'table': self._create_table_slide,
            'quote': self._create_quote_slide,
            'conclusion': self._create_conclusion_slide,
            'metrics': self._create_metrics_slide,
            'timeline': self._create_timeline_slide,
            'comparison': self._create_comparison_slide,
            'image': self._create_image_slide
        }
        
        creator = layout_creators.get(layout_type, self._create_content_slide)
        creator(slide, slide_content)
        
        self._add_slide_number(slide, slide_content.get('slide_number', 0))
    
    def _add_background(self, slide):
        background = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE, 0, 0,
            Inches(self.template_data['slide_size']['width']),
            Inches(self.template_data['slide_size']['height'])
        )
        background.fill.solid()
        background.fill.fore_color.rgb = self._hex_to_rgb(self.template_data['colors']['background'])
        background.line.fill.background()
        
        spTree = slide.shapes._spTree
        sp = background._element
        spTree.remove(sp)
        spTree.insert(2, sp)
    
    def _create_title_slide(self, slide, content: Dict):
        title_box = slide.shapes.add_textbox(Inches(0.5), Inches(2.5), Inches(12.333), Inches(1.5))
        frame = title_box.text_frame
        para = frame.paragraphs[0]
        para.text = content.get('title', 'Presentation Title')
        para.font.size = Pt(self.template_data['fonts']['title']['size'])
        para.font.bold = True
        para.font.color.rgb = self._hex_to_rgb(self.template_data['colors']['text_primary'])
        para.alignment = PP_ALIGN.CENTER
        
        if content.get('subtitle'):
            subtitle_box = slide.shapes.add_textbox(Inches(0.5), Inches(4.2), Inches(12.333), Inches(1))
            frame = subtitle_box.text_frame
            para = frame.paragraphs[0]
            para.text = content.get('subtitle', '')
            para.font.size = Pt(self.template_data['fonts']['subtitle']['size'])
            para.font.color.rgb = self._hex_to_rgb(self.template_data['colors']['text_secondary'])
            para.alignment = PP_ALIGN.CENTER
        
        # Decorative line
        line = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(5), Inches(4), Inches(3.333), Inches(0.05))
        line.fill.solid()
        line.fill.fore_color.rgb = self._hex_to_rgb(self.template_data['colors']['primary'])
        line.line.fill.background()
    
    def _create_content_slide(self, slide, content: Dict):
        self._add_slide_title(slide, content.get('title', ''))
        
        content_data = content.get('content', {})
        if content_data.get('bullet_points'):
            self._add_bullet_points(slide, content_data['bullet_points'], 
                                   Inches(0.75), Inches(1.8), Inches(11.833), Inches(5))
        elif content_data.get('main_text'):
            self._add_text_box(slide, content_data['main_text'],
                              Inches(0.75), Inches(1.8), Inches(11.833), Inches(5))
    
    def _create_two_column_slide(self, slide, content: Dict):
        self._add_slide_title(slide, content.get('title', ''))
        
        content_data = content.get('content', {})
        
        left = content_data.get('left_column', '')
        if isinstance(left, list):
            self._add_bullet_points(slide, left, Inches(0.5), Inches(1.8), Inches(5.8), Inches(5))
        else:
            self._add_text_box(slide, str(left), Inches(0.5), Inches(1.8), Inches(5.8), Inches(5))
        
        line = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(6.5), Inches(1.8), Inches(0.02), Inches(4.5))
        line.fill.solid()
        line.fill.fore_color.rgb = self._hex_to_rgb(self.template_data['colors']['secondary'])
        line.line.fill.background()
        
        right = content_data.get('right_column', '')
        if isinstance(right, list):
            self._add_bullet_points(slide, right, Inches(7), Inches(1.8), Inches(5.8), Inches(5))
        else:
            self._add_text_box(slide, str(right), Inches(7), Inches(1.8), Inches(5.8), Inches(5))
    
    def _create_chart_slide(self, slide, content: Dict):
        self._add_slide_title(slide, content.get('title', ''))
        
        content_data = content.get('content', {})
        chart_info = content_data.get('chart', {})
        
        if chart_info:
            chart_type = chart_info.get('type', 'bar')
            chart_data = {
                'title': chart_info.get('title', ''),
                'labels': chart_info.get('labels', ['A', 'B', 'C', 'D']),
                'datasets': chart_info.get('datasets', [{'name': 'Series 1', 'values': [65, 78, 90, 85]}]),
                'x_axis_label': chart_info.get('x_axis_label', ''),
                'y_axis_label': chart_info.get('y_axis_label', '')
            }
            
            try:
                chart_image = self.chart_generator.create_chart(chart_type, chart_data, 900, 450)
                image_stream = io.BytesIO(chart_image)
                slide.shapes.add_picture(image_stream, Inches(1.5), Inches(1.8), Inches(10), Inches(5))
            except:
                self._add_text_box(slide, f"[Chart: {chart_info.get('title', 'Data')}]",
                                  Inches(1.5), Inches(3), Inches(10), Inches(2))
    
    def _create_table_slide(self, slide, content: Dict):
        self._add_slide_title(slide, content.get('title', ''))
        
        content_data = content.get('content', {})
        table_info = content_data.get('table', {})
        
        headers = table_info.get('headers', ['Column 1', 'Column 2'])
        rows = table_info.get('rows', [['Data 1', 'Data 2']])
        
        num_cols = len(headers)
        num_rows = len(rows) + 1
        
        table = slide.shapes.add_table(num_rows, num_cols, Inches(1.5), Inches(2), Inches(10), Inches(min(num_rows * 0.6, 4.5))).table
        
        for i, header in enumerate(headers):
            cell = table.cell(0, i)
            cell.text = str(header)
            cell.fill.solid()
            cell.fill.fore_color.rgb = self._hex_to_rgb(self.template_data['colors']['primary'])
            para = cell.text_frame.paragraphs[0]
            para.font.bold = True
            para.font.color.rgb = RGBColor(255, 255, 255)
            para.font.size = Pt(14)
            para.alignment = PP_ALIGN.CENTER
        
        for row_idx, row_data in enumerate(rows):
            for col_idx, cell_data in enumerate(row_data):
                if col_idx < num_cols:
                    cell = table.cell(row_idx + 1, col_idx)
                    cell.text = str(cell_data)
                    cell.fill.solid()
                    cell.fill.fore_color.rgb = self._hex_to_rgb('#1E293B' if row_idx % 2 == 0 else '#0F172A')
                    para = cell.text_frame.paragraphs[0]
                    para.font.color.rgb = self._hex_to_rgb(self.template_data['colors']['text_primary'])
                    para.font.size = Pt(12)
                    para.alignment = PP_ALIGN.CENTER
    
    def _create_quote_slide(self, slide, content: Dict):
        content_data = content.get('content', {})
        quote = content_data.get('quote', 'Add your quote here')
        author = content_data.get('quote_author', '')
        
        quote_mark = slide.shapes.add_textbox(Inches(1), Inches(2), Inches(1), Inches(1))
        frame = quote_mark.text_frame
        para = frame.paragraphs[0]
        para.text = '"'
        para.font.size = Pt(120)
        para.font.color.rgb = self._hex_to_rgb(self.template_data['colors']['primary'])
        
        quote_box = slide.shapes.add_textbox(Inches(2), Inches(2.5), Inches(9), Inches(3))
        frame = quote_box.text_frame
        frame.word_wrap = True
        para = frame.paragraphs[0]
        para.text = quote
        para.font.size = Pt(28)
        para.font.italic = True
        para.font.color.rgb = self._hex_to_rgb(self.template_data['colors']['text_primary'])
        para.alignment = PP_ALIGN.CENTER
        
        if author:
            author_box = slide.shapes.add_textbox(Inches(2), Inches(5.5), Inches(9), Inches(0.5))
            frame = author_box.text_frame
            para = frame.paragraphs[0]
            para.text = f"— {author}"
            para.font.size = Pt(20)
            para.font.color.rgb = self._hex_to_rgb(self.template_data['colors']['secondary'])
            para.alignment = PP_ALIGN.RIGHT
    
    def _create_conclusion_slide(self, slide, content: Dict):
        title_box = slide.shapes.add_textbox(Inches(0.5), Inches(2), Inches(12.333), Inches(1.5))
        frame = title_box.text_frame
        para = frame.paragraphs[0]
        para.text = content.get('title', 'Thank You')
        para.font.size = Pt(48)
        para.font.bold = True
        para.font.color.rgb = self._hex_to_rgb(self.template_data['colors']['text_primary'])
        para.alignment = PP_ALIGN.CENTER
        
        content_data = content.get('content', {})
        bullet_points = content_data.get('bullet_points', [])
        
        if bullet_points:
            self._add_bullet_points(slide, bullet_points, Inches(2), Inches(3.5), Inches(9.333), Inches(3), centered=True)
        
        cta = content_data.get('call_to_action', '')
        if cta:
            cta_box = slide.shapes.add_textbox(Inches(0.5), Inches(6.5), Inches(12.333), Inches(0.5))
            frame = cta_box.text_frame
            para = frame.paragraphs[0]
            para.text = cta
            para.font.size = Pt(18)
            para.font.color.rgb = self._hex_to_rgb(self.template_data['colors']['accent'])
            para.alignment = PP_ALIGN.CENTER
    
    def _create_metrics_slide(self, slide, content: Dict):
        self._add_slide_title(slide, content.get('title', 'Key Metrics'))
        
        content_data = content.get('content', {})
        metrics = content_data.get('key_metrics', [])
        
        if metrics:
            num_metrics = min(len(metrics), 4)
            metric_width = 10 / num_metrics
            
            for i, metric in enumerate(metrics[:4]):
                x_pos = 1.5 + (i * metric_width)
                
                box = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(x_pos), Inches(2.5), Inches(metric_width - 0.3), Inches(3))
                box.fill.solid()
                box.fill.fore_color.rgb = self._hex_to_rgb('#1E293B')
                box.line.color.rgb = self._hex_to_rgb(self.template_data['colors']['primary'])
                
                value_box = slide.shapes.add_textbox(Inches(x_pos), Inches(3), Inches(metric_width - 0.3), Inches(1))
                frame = value_box.text_frame
                para = frame.paragraphs[0]
                para.text = str(metric.get('value', '0'))
                para.font.size = Pt(36)
                para.font.bold = True
                para.font.color.rgb = self._hex_to_rgb(self.template_data['colors']['accent'])
                para.alignment = PP_ALIGN.CENTER
                
                label_box = slide.shapes.add_textbox(Inches(x_pos), Inches(4.2), Inches(metric_width - 0.3), Inches(0.8))
                frame = label_box.text_frame
                frame.word_wrap = True
                para = frame.paragraphs[0]
                para.text = metric.get('label', 'Metric')
                para.font.size = Pt(16)
                para.font.color.rgb = self._hex_to_rgb(self.template_data['colors']['text_secondary'])
                para.alignment = PP_ALIGN.CENTER
    
    def _create_timeline_slide(self, slide, content: Dict):
        self._add_slide_title(slide, content.get('title', 'Timeline'))
        
        content_data = content.get('content', {})
        items = content_data.get('timeline_items', [])
        
        if items:
            line = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(1), Inches(4), Inches(11.333), Inches(0.05))
            line.fill.solid()
            line.fill.fore_color.rgb = self._hex_to_rgb(self.template_data['colors']['primary'])
            line.line.fill.background()
            
            item_width = 11.333 / len(items[:6])
            for i, item in enumerate(items[:6]):
                x_pos = 1 + (i * item_width)
                
                circle = slide.shapes.add_shape(MSO_SHAPE.OVAL, Inches(x_pos + item_width/2 - 0.15), Inches(3.85), Inches(0.3), Inches(0.3))
                circle.fill.solid()
                circle.fill.fore_color.rgb = self._hex_to_rgb(self.template_data['colors']['accent'])
                circle.line.fill.background()
                
                date_box = slide.shapes.add_textbox(Inches(x_pos), Inches(2.5), Inches(item_width), Inches(0.5))
                frame = date_box.text_frame
                para = frame.paragraphs[0]
                para.text = item.get('date', '')
                para.font.size = Pt(14)
                para.font.bold = True
                para.font.color.rgb = self._hex_to_rgb(self.template_data['colors']['secondary'])
                para.alignment = PP_ALIGN.CENTER
                
                event_box = slide.shapes.add_textbox(Inches(x_pos), Inches(4.3), Inches(item_width), Inches(1.5))
                frame = event_box.text_frame
                frame.word_wrap = True
                para = frame.paragraphs[0]
                para.text = item.get('event', '')
                para.font.size = Pt(12)
                para.font.color.rgb = self._hex_to_rgb(self.template_data['colors']['text_primary'])
                para.alignment = PP_ALIGN.CENTER
    
    def _create_comparison_slide(self, slide, content: Dict):
        self._add_slide_title(slide, content.get('title', 'Comparison'))
        
        content_data = content.get('content', {})
        items = content_data.get('comparison_items', [])
        
        if len(items) >= 2:
            for idx, item in enumerate(items[:2]):
                x_start = Inches(0.5) if idx == 0 else Inches(7)
                
                box = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x_start, Inches(1.8), Inches(5.8), Inches(5))
                box.fill.solid()
                box.fill.fore_color.rgb = self._hex_to_rgb('#1E293B')
                box.line.color.rgb = self._hex_to_rgb(self.template_data['colors']['primary'])
                
                title_box = slide.shapes.add_textbox(x_start + Inches(0.3), Inches(2), Inches(5.2), Inches(0.6))
                frame = title_box.text_frame
                para = frame.paragraphs[0]
                para.text = item.get('title', f'Option {idx + 1}')
                para.font.size = Pt(24)
                para.font.bold = True
                para.font.color.rgb = self._hex_to_rgb(self.template_data['colors']['text_primary'])
                para.alignment = PP_ALIGN.CENTER
            
            vs_box = slide.shapes.add_textbox(Inches(6), Inches(3.5), Inches(1.333), Inches(1))
            frame = vs_box.text_frame
            para = frame.paragraphs[0]
            para.text = "VS"
            para.font.size = Pt(36)
            para.font.bold = True
            para.font.color.rgb = self._hex_to_rgb(self.template_data['colors']['accent'])
            para.alignment = PP_ALIGN.CENTER
    
    def _create_image_slide(self, slide, content: Dict):
        self._add_slide_title(slide, content.get('title', ''))
        
        content_data = content.get('content', {})
        image_desc = content_data.get('image_description', 'Image Placeholder')
        
        placeholder = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(1.5), Inches(1.8), Inches(10), Inches(5))
        placeholder.fill.solid()
        placeholder.fill.fore_color.rgb = self._hex_to_rgb('#1E293B')
        placeholder.line.color.rgb = self._hex_to_rgb(self.template_data['colors']['secondary'])
        placeholder.line.width = Pt(2)
        
        text_box = slide.shapes.add_textbox(Inches(1.5), Inches(4), Inches(10), Inches(1))
        frame = text_box.text_frame
        para = frame.paragraphs[0]
        para.text = f"📷 {image_desc}"
        para.font.size = Pt(18)
        para.font.color.rgb = self._hex_to_rgb(self.template_data['colors']['text_secondary'])
        para.alignment = PP_ALIGN.CENTER
    
    def _add_slide_title(self, slide, title: str):
        title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.4), Inches(12.333), Inches(1))
        frame = title_box.text_frame
        para = frame.paragraphs[0]
        para.text = title
        para.font.size = Pt(self.template_data['fonts']['title']['size'])
        para.font.bold = True
        para.font.color.rgb = self._hex_to_rgb(self.template_data['colors']['text_primary'])
        
        line = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.5), Inches(1.4), Inches(2), Inches(0.04))
        line.fill.solid()
        line.fill.fore_color.rgb = self._hex_to_rgb(self.template_data['colors']['primary'])
        line.line.fill.background()
    
    def _add_bullet_points(self, slide, points: List[str], left, top, width, height, centered: bool = False):
        text_box = slide.shapes.add_textbox(left, top, width, height)
        frame = text_box.text_frame
        frame.word_wrap = True
        
        for i, point in enumerate(points):
            if i == 0:
                para = frame.paragraphs[0]
            else:
                para = frame.add_paragraph()
            
            para.text = f"• {point}"
            para.font.size = Pt(self.template_data['fonts']['body']['size'])
            para.font.color.rgb = self._hex_to_rgb(self.template_data['colors']['text_primary'])
            para.space_after = Pt(12)
            
            if centered:
                para.alignment = PP_ALIGN.CENTER
    
    def _add_text_box(self, slide, text: str, left, top, width, height):
        text_box = slide.shapes.add_textbox(left, top, width, height)
        frame = text_box.text_frame
        frame.word_wrap = True
        
        para = frame.paragraphs[0]
        para.text = text
        para.font.size = Pt(self.template_data['fonts']['body']['size'])
        para.font.color.rgb = self._hex_to_rgb(self.template_data['colors']['text_primary'])
    
    def _add_slide_number(self, slide, number: int):
        if number <= 0:
            return
        
        num_box = slide.shapes.add_textbox(Inches(12), Inches(7), Inches(1), Inches(0.4))
        frame = num_box.text_frame
        para = frame.paragraphs[0]
        para.text = str(number)
        para.font.size = Pt(12)
        para.font.color.rgb = self._hex_to_rgb(self.template_data['colors']['text_secondary'])
        para.alignment = PP_ALIGN.RIGHT
    
    def _hex_to_rgb(self, hex_color: str) -> RGBColor:
        hex_color = hex_color.lstrip('#')
        return RGBColor(int(hex_color[0:2], 16), int(hex_color[2:4], 16), int(hex_color[4:6], 16))
