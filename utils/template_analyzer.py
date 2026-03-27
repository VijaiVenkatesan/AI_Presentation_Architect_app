"""
Template Analyzer Module - Dynamic Template Extraction
Handles any PowerPoint template type and extracts ALL styling
"""

import io
import re
import copy
from typing import Dict, List, Any, Optional, Tuple
from PIL import Image
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE_TYPE
from collections import Counter
import colorsys


class TemplateAnalyzer:
    """Analyzes any PowerPoint template and extracts complete styling"""
    
    def __init__(self):
        self.template_data = self._get_default_template()
        self.original_prs = None
        self.extracted_images = {}
    
    def _get_default_template(self) -> Dict:
        """Default template settings"""
        return {
            'colors': {
                'primary': '#6366F1',
                'secondary': '#8B5CF6',
                'accent': '#EC4899',
                'background': '#FFFFFF',
                'text_primary': '#1F2937',
                'text_secondary': '#6B7280'
            },
            'fonts': {
                'title': {'name': 'Arial', 'size': 44, 'bold': True, 'color': '#1F2937'},
                'subtitle': {'name': 'Arial', 'size': 28, 'bold': False, 'color': '#6B7280'},
                'body': {'name': 'Arial', 'size': 18, 'bold': False, 'color': '#1F2937'}
            },
            'layouts': [],
            'slide_size': {'width': 13.333, 'height': 7.5},
            'has_logo': False,
            'logo_image': None,
            'logo_position': {'left': 0.3, 'top': 0.3, 'width': 1.5, 'height': 0.75},
            'background': {'type': 'solid', 'color': '#FFFFFF'},
            'color_scheme': [],
            'master_layouts': [],
            'use_template_file': False,
            'template_bytes': None,
            'slide_layouts_info': []
        }
    
    def analyze_pptx(self, pptx_file: io.BytesIO) -> Dict[str, Any]:
        """Analyze PowerPoint and extract ALL template information"""
        try:
            # Store original bytes
            pptx_file.seek(0)
            self.template_data['template_bytes'] = pptx_file.read()
            pptx_file.seek(0)
            self.template_data['use_template_file'] = True
            
            # Parse presentation
            prs = Presentation(pptx_file)
            self.original_prs = prs
            
            # Extract slide size
            self.template_data['slide_size'] = {
                'width': prs.slide_width.inches,
                'height': prs.slide_height.inches
            }
            
            # Extract from slide masters FIRST (most important)
            self._extract_from_masters(prs)
            
            # Extract from actual slides
            all_colors = []
            all_fonts = []
            
            for slide_idx, slide in enumerate(prs.slides):
                try:
                    # Extract background
                    self._extract_slide_background(slide)
                    
                    # Extract shapes
                    for shape in slide.shapes:
                        try:
                            # Extract colors from shape
                            shape_colors = self._extract_shape_colors(shape)
                            all_colors.extend(shape_colors)
                            
                            # Extract fonts from shape
                            shape_fonts = self._extract_shape_fonts(shape)
                            all_fonts.extend(shape_fonts)
                            
                            # Extract logo (images in corners)
                            if slide_idx == 0:
                                self._check_for_logo(shape)
                        except Exception:
                            continue
                except Exception as e:
                    print(f"Error on slide {slide_idx}: {e}")
                    continue
            
            # Process all extracted data
            self._analyze_colors(all_colors)
            self._analyze_fonts(all_fonts)
            
            # Store layout information for each slide layout
            self._extract_layout_info(prs)
            
            return self.template_data
            
        except Exception as e:
            print(f"Error analyzing PPTX: {e}")
            import traceback
            traceback.print_exc()
            return self.template_data
    
    def _extract_from_masters(self, prs: Presentation):
        """Extract styles from slide masters"""
        try:
            for master in prs.slide_masters:
                # Extract background from master
                try:
                    if hasattr(master, 'background') and master.background:
                        self._extract_background_fill(master.background)
                except:
                    pass
                
                # Extract shapes from master (logos, decorations)
                try:
                    for shape in master.shapes:
                        # Check for logo
                        self._check_for_logo(shape)
                        
                        # Extract colors
                        colors = self._extract_shape_colors(shape)
                        if colors:
                            self.template_data['color_scheme'].extend(colors)
                except:
                    pass
                
                # Extract from slide layouts
                try:
                    for layout in master.slide_layouts:
                        layout_info = {
                            'name': layout.name,
                            'shapes': []
                        }
                        
                        for shape in layout.shapes:
                            try:
                                if hasattr(shape, 'is_placeholder') and shape.is_placeholder:
                                    ph = shape.placeholder_format
                                    shape_data = {
                                        'placeholder_type': str(ph.type) if ph else 'unknown',
                                        'left': shape.left.inches if shape.left else 0,
                                        'top': shape.top.inches if shape.top else 0,
                                        'width': shape.width.inches if shape.width else 0,
                                        'height': shape.height.inches if shape.height else 0
                                    }
                                    
                                    # Extract font from placeholder
                                    fonts = self._extract_shape_fonts(shape)
                                    if fonts:
                                        shape_data['font'] = fonts[0]
                                    
                                    layout_info['shapes'].append(shape_data)
                            except:
                                continue
                        
                        self.template_data['master_layouts'].append(layout_info)
                except:
                    pass
                    
        except Exception as e:
            print(f"Error extracting from masters: {e}")
    
    def _extract_slide_background(self, slide):
        """Extract background from slide"""
        try:
            if hasattr(slide, 'background') and slide.background:
                self._extract_background_fill(slide.background)
        except:
            pass
    
    def _extract_background_fill(self, background):
        """Extract fill information from background"""
        try:
            fill = background.fill
            
            if fill is None or fill.type is None:
                return
            
            fill_type = str(fill.type)
            
            if 'SOLID' in fill_type:
                try:
                    if fill.fore_color and fill.fore_color.rgb:
                        color = self._rgb_to_hex(fill.fore_color.rgb)
                        self.template_data['background'] = {
                            'type': 'solid',
                            'color': color
                        }
                        self.template_data['colors']['background'] = color
                except:
                    pass
            
            elif 'GRADIENT' in fill_type:
                self.template_data['background']['type'] = 'gradient'
                try:
                    gradient_colors = []
                    for stop in fill.gradient_stops:
                        if stop.color and stop.color.rgb:
                            gradient_colors.append(self._rgb_to_hex(stop.color.rgb))
                    if gradient_colors:
                        self.template_data['background']['gradient_colors'] = gradient_colors
                        self.template_data['colors']['background'] = gradient_colors[0]
                except:
                    pass
                    
        except Exception as e:
            print(f"Error extracting background: {e}")
    
    def _extract_shape_colors(self, shape) -> List[str]:
        """Extract all colors from a shape"""
        colors = []
        
        try:
            # Fill color
            if hasattr(shape, 'fill') and shape.fill:
                try:
                    if shape.fill.type is not None and 'SOLID' in str(shape.fill.type):
                        if shape.fill.fore_color and shape.fill.fore_color.rgb:
                            colors.append(self._rgb_to_hex(shape.fill.fore_color.rgb))
                except:
                    pass
            
            # Line color
            if hasattr(shape, 'line') and shape.line:
                try:
                    if shape.line.color and shape.line.color.rgb:
                        colors.append(self._rgb_to_hex(shape.line.color.rgb))
                except:
                    pass
            
            # Text colors
            if hasattr(shape, 'text_frame') and hasattr(shape, 'has_text_frame') and shape.has_text_frame:
                try:
                    for para in shape.text_frame.paragraphs:
                        for run in para.runs:
                            if run.font.color and run.font.color.type is not None:
                                if run.font.color.rgb:
                                    colors.append(self._rgb_to_hex(run.font.color.rgb))
                except:
                    pass
        except:
            pass
        
        return colors
    
    def _extract_shape_fonts(self, shape) -> List[Dict]:
        """Extract font information from shape"""
        fonts = []
        
        try:
            if not hasattr(shape, 'text_frame') or not hasattr(shape, 'has_text_frame'):
                return fonts
            
            if not shape.has_text_frame:
                return fonts
            
            for para in shape.text_frame.paragraphs:
                for run in para.runs:
                    try:
                        font_info = {
                            'name': run.font.name if run.font.name else 'Arial',
                            'size': run.font.size.pt if run.font.size else 18,
                            'bold': run.font.bold if run.font.bold is not None else False,
                            'italic': run.font.italic if run.font.italic is not None else False,
                            'color': '#1F2937'
                        }
                        
                        # Get font color
                        try:
                            if run.font.color and run.font.color.type is not None:
                                if run.font.color.rgb:
                                    font_info['color'] = self._rgb_to_hex(run.font.color.rgb)
                        except:
                            pass
                        
                        if font_info['name'] and font_info['size']:
                            fonts.append(font_info)
                    except:
                        continue
        except:
            pass
        
        return fonts
    
    def _check_for_logo(self, shape):
        """Check if shape is a logo and extract it"""
        try:
            if shape.shape_type != MSO_SHAPE_TYPE.PICTURE:
                return
            
            left = shape.left.inches if shape.left else 0
            top = shape.top.inches if shape.top else 0
            width = shape.width.inches if shape.width else 0
            height = shape.height.inches if shape.height else 0
            
            slide_width = self.template_data['slide_size']['width']
            slide_height = self.template_data['slide_size']['height']
            
            # Logo criteria: small image in any corner
            is_top = top < 1.5
            is_bottom = top > slide_height - 2
            is_left = left < 2
            is_right = left > slide_width - 3
            is_corner = (is_top or is_bottom) and (is_left or is_right)
            is_small = width < 3 and height < 2
            
            if is_corner and is_small:
                try:
                    if hasattr(shape, 'image') and shape.image:
                        self.template_data['has_logo'] = True
                        self.template_data['logo_image'] = shape.image.blob
                        self.template_data['logo_position'] = {
                            'left': left,
                            'top': top,
                            'width': width,
                            'height': height
                        }
                except:
                    pass
        except:
            pass
    
    def _extract_layout_info(self, prs: Presentation):
        """Extract information about available slide layouts"""
        try:
            layouts_info = []
            
            for layout in prs.slide_layouts:
                layout_data = {
                    'name': layout.name,
                    'index': prs.slide_layouts.index(layout)
                }
                layouts_info.append(layout_data)
            
            self.template_data['slide_layouts_info'] = layouts_info
        except:
            pass
    
    def _analyze_colors(self, colors: List[str]):
        """Analyze extracted colors and categorize them"""
        if not colors:
            return
        
        # Filter valid colors
        valid_colors = []
        for c in colors:
            if c and isinstance(c, str) and re.match(r'^#[0-9A-Fa-f]{6}$', c):
                valid_colors.append(c)
        
        if not valid_colors:
            return
        
        # Count occurrences
        color_counts = Counter(valid_colors)
        sorted_colors = [color for color, _ in color_counts.most_common(15)]
        
        # Categorize by lightness
        dark_colors = []
        light_colors = []
        mid_colors = []
        
        for color in sorted_colors:
            lightness = self._get_lightness(color)
            if lightness < 0.25:
                dark_colors.append(color)
            elif lightness > 0.75:
                light_colors.append(color)
            else:
                mid_colors.append(color)
        
        # Assign to color scheme
        # Background: darkest or lightest depending on overall theme
        if dark_colors and len(dark_colors) >= len(light_colors):
            # Dark theme
            self.template_data['colors']['background'] = dark_colors[0]
            if light_colors:
                self.template_data['colors']['text_primary'] = light_colors[0]
            if len(light_colors) > 1:
                self.template_data['colors']['text_secondary'] = light_colors[1]
        elif light_colors:
            # Light theme
            self.template_data['colors']['background'] = light_colors[0]
            if dark_colors:
                self.template_data['colors']['text_primary'] = dark_colors[0]
            if len(dark_colors) > 1:
                self.template_data['colors']['text_secondary'] = dark_colors[1]
        
        # Accent colors from mid-range
        if mid_colors:
            self.template_data['colors']['primary'] = mid_colors[0]
            if len(mid_colors) > 1:
                self.template_data['colors']['secondary'] = mid_colors[1]
            if len(mid_colors) > 2:
                self.template_data['colors']['accent'] = mid_colors[2]
        
        # Store full color scheme
        self.template_data['color_scheme'] = sorted_colors[:10]
    
    def _analyze_fonts(self, fonts: List[Dict]):
        """Analyze extracted fonts and categorize them"""
        if not fonts:
            return
        
        # Filter fonts with valid size
        sized_fonts = [(f, f.get('size', 0)) for f in fonts if f.get('size') and f.get('size') > 0]
        
        if not sized_fonts:
            return
        
        # Sort by size (largest first)
        sized_fonts.sort(key=lambda x: x[1], reverse=True)
        
        # Get unique font names with their typical sizes
        font_by_size = {}
        for font, size in sized_fonts:
            if size not in font_by_size:
                font_by_size[size] = font
        
        sizes = sorted(font_by_size.keys(), reverse=True)
        
        # Assign fonts based on size
        if sizes:
            # Largest = title
            largest_size = sizes[0]
            title_font = font_by_size[largest_size]
            self.template_data['fonts']['title'] = {
                'name': title_font.get('name', 'Arial'),
                'size': int(largest_size),
                'bold': title_font.get('bold', True),
                'color': title_font.get('color', self.template_data['colors']['text_primary'])
            }
            
            # Medium = subtitle
            if len(sizes) > 1:
                mid_size = sizes[len(sizes)//3] if len(sizes) > 2 else sizes[1]
                mid_font = font_by_size[mid_size]
                self.template_data['fonts']['subtitle'] = {
                    'name': mid_font.get('name', 'Arial'),
                    'size': int(mid_size),
                    'bold': mid_font.get('bold', False),
                    'color': mid_font.get('color', self.template_data['colors']['text_secondary'])
                }
            
            # Smallest = body
            if len(sizes) > 2:
                small_size = sizes[-1]
                small_font = font_by_size[small_size]
                self.template_data['fonts']['body'] = {
                    'name': small_font.get('name', 'Arial'),
                    'size': max(int(small_size), 12),  # Minimum 12pt
                    'bold': small_font.get('bold', False),
                    'color': small_font.get('color', self.template_data['colors']['text_primary'])
                }
    
    def analyze_image(self, image_file: io.BytesIO) -> Dict[str, Any]:
        """Analyze image to extract colors"""
        try:
            image = Image.open(image_file)
            colors = self._extract_colors_from_image(image)
            self._analyze_colors(colors)
            return self.template_data
        except Exception as e:
            print(f"Error analyzing image: {e}")
            return self.template_data
    
    def _extract_colors_from_image(self, image: Image.Image) -> List[str]:
        """Extract colors from image"""
        try:
            image = image.resize((100, 100))
            if image.mode != 'RGB':
                image = image.convert('RGB')
            
            pixels = list(image.getdata())
            color_counts = Counter()
            
            for pixel in pixels:
                # Round to reduce unique colors
                rounded = (pixel[0] // 16 * 16, pixel[1] // 16 * 16, pixel[2] // 16 * 16)
                color_counts[rounded] += 1
            
            top_colors = color_counts.most_common(15)
            return [f'#{r:02X}{g:02X}{b:02X}' for (r, g, b), _ in top_colors]
        except:
            return []
    
    def _rgb_to_hex(self, rgb) -> str:
        """Convert RGB to hex"""
        try:
            if hasattr(rgb, '__iter__') and len(rgb) >= 3:
                return f'#{rgb[0]:02X}{rgb[1]:02X}{rgb[2]:02X}'
        except:
            pass
        return '#000000'
    
    def _get_lightness(self, hex_color: str) -> float:
        """Get lightness of color (0-1)"""
        try:
            hex_color = hex_color.lstrip('#')
            r = int(hex_color[0:2], 16) / 255
            g = int(hex_color[2:4], 16) / 255
            b = int(hex_color[4:6], 16) / 255
            _, lightness, _ = colorsys.rgb_to_hls(r, g, b)
            return lightness
        except:
            return 0.5
    
    def get_template_summary(self) -> str:
        """Generate summary"""
        lines = ["## Template Analysis\n"]
        
        lines.append("### Colors")
        for name, color in self.template_data['colors'].items():
            lines.append(f"- {name}: {color}")
        
        lines.append("\n### Fonts")
        for name, font in self.template_data['fonts'].items():
            lines.append(f"- {name}: {font.get('name')} {font.get('size')}pt")
        
        lines.append(f"\n### Logo: {'✓ Found' if self.template_data['has_logo'] else '✗ Not found'}")
        lines.append(f"### Background: {self.template_data['background'].get('type', 'solid')}")
        
        return "\n".join(lines)
