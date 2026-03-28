from pptx import Presentation
from pptx.util import Inches
import json
import os

from core.clone_engine import clone_slide, replace_text
from core.image_engine import generate_image
from core.diagram_engine import add_diagram


def generate_ppt(template_file, content_json, output_file="output.pptx"):

    # Load template
    prs = Presentation(template_file)

    # Parse JSON safely
    try:
        data = json.loads(content_json)
    except Exception as e:
        raise Exception(f"Invalid JSON content: {e}")

    slides_data = data.get("slides", [])

    if not slides_data:
        raise Exception("No slides data found")

    base_slide_index = 0  # Use first slide as design base

    for idx, slide_data in enumerate(slides_data):

        try:
            # Clone template slide (pixel-level)
            slide = clone_slide(prs, base_slide_index)

            # Replace text content
            replace_text(slide, slide_data)

            # -------------------------
            # 🖼 IMAGE INSERTION
            # -------------------------
            img_prompt = slide_data.get("image_prompt")

            if img_prompt:
                try:
                    img = generate_image(img_prompt)
                    if img:
                        img_path = f"temp_{idx}.png"
                        img.save(img_path)

                        slide.shapes.add_picture(
                            img_path,
                            Inches(1),
                            Inches(4),
                            width=Inches(4)
                        )

                        os.remove(img_path)

                except Exception as e:
                    print(f"Image generation failed: {e}")

            # -------------------------
            # 📊 DIAGRAM INSERTION
            # -------------------------
            diagram_type = slide_data.get("diagram_type")

            if diagram_type:
                try:
                    add_diagram(slide, diagram_type)
                except Exception as e:
                    print(f"Diagram error: {e}")

        except Exception as e:
            print(f"Slide {idx} failed: {e}")
            continue  # Skip bad slide but continue

    # Save final PPT
    prs.save(output_file)

    return output_file
